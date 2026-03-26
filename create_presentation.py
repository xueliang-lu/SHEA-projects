#!/usr/bin/env python3
"""
SHEA Projects Presentation Generator
Creates a comprehensive PPT for all 3 SHEA projects
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # Set slide size to 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    def add_title_slide(title, subtitle=""):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_shape = slide.shapes.title
        title_shape.text = title
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
        return slide
    
    def add_content_slide(title, content_items):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        title_shape.text = title
        
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        
        for i, item in enumerate(content_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
            p.space_after = Pt(10)
        
        return slide
    
    def add_two_column_slide(title, left_title, left_content, right_title, right_content):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Left column
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(6), Inches(5))
        left_tf = left_box.text_frame
        left_tf.word_wrap = True
        
        p = left_tf.paragraphs[0]
        p.text = left_title
        p.font.bold = True
        p.font.size = Pt(20)
        
        for item in left_content:
            p = left_tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(16)
            p.space_after = Pt(8)
        
        # Right column
        right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(6), Inches(5))
        right_tf = right_box.text_frame
        right_tf.word_wrap = True
        
        p = right_tf.paragraphs[0]
        p.text = right_title
        p.font.bold = True
        p.font.size = Pt(20)
        
        for item in right_content:
            p = right_tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(16)
            p.space_after = Pt(8)
        
        return slide
    
    # ========== SLIDE 1: Title Slide ==========
    add_title_slide(
        "SHEA Projects Portfolio",
        "Sydney Higher Education Academy\nInnovation & Automation Systems\n\nPresented by: Sunil Paudel & Xueliang Lu"
    )
    
    # ========== SLIDE 2: Overview ==========
    add_content_slide(
        "Project Overview",
        [
            "Three integrated systems for SHEA operations:",
            "",
            "1. SHEAC-HATBOT — AI Chatbot with RAG",
            "   • Student & staff support automation",
            "   • FAQ knowledge base with AI responses",
            "",
            "2. CPL Automation — Credit for Prior Learning Assistant",
            "   • Automated unit mapping & confidence scoring",
            "   • Streamlines credit transfer process",
            "",
            "3. Student Risk Analytics Dashboard",
            "   • Early intervention system",
            "   • Moodle integration with risk scoring"
        ]
    )
    
    # ========== SLIDE 3: Project 1 - Title ==========
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    title_shape.text = "Project 1: SHEAC-HATBOT"
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = "AI Chatbot with RAG System"
    
    # ========== SLIDE 4: HATBOT Features ==========
    add_two_column_slide(
        "SHEAC-HATBOT — Key Features",
        "Core Capabilities",
        [
            "Dual AI Engines (OpenAI GPT-4o-mini & Google Gemini 2.0 Flash)",
            "RAG system for FAQ knowledge base",
            "Context-aware, accurate answers",
            "Full chat history with search",
            "Persistent MongoDB storage",
            "Modern, responsive UI"
        ],
        "Technical Stack",
        [
            "Next.js 16.1.6 (App Router)",
            "TypeScript 5",
            "Tailwind CSS 4",
            "React 19",
            "MongoDB (Mongoose)",
            "OpenAI & Google Gemini APIs"
        ]
    )
    
    # ========== SLIDE 5: HATBOT Architecture ==========
    add_content_slide(
        "SHEAC-HATBOT — How It Works",
        [
            "RAG (Retrieval Augmented Generation) Flow:",
            "",
            "1. User asks a question",
            "2. System searches FAQ documents in data/ folder",
            "3. Retrieves relevant passages",
            "4. Augments AI prompt with context",
            "5. Returns accurate, context-aware answer",
            "",
            "Key Benefits:",
            "• Reduces AI hallucinations",
            "• Provides institution-specific answers",
            "• Easy to update knowledge base (just add files to data/)",
            "• Switch between AI providers seamlessly"
        ]
    )
    
    # ========== SLIDE 6: Project 2 - Title ==========
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    title_shape.text = "Project 2: CPL Automation"
    subtitle_shape = slide.placeholders[1]
    subtitle_shape = "Credit for Prior Learning Assistant"
    
    # ========== SLIDE 7: CPL Features ==========
    add_two_column_slide(
        "CPL Automation — Key Features",
        "What It Does",
        [
            "Maps external university units to SHEA units",
            "Confidence scoring with explanations",
            "Review & approval workflow",
            "Exportable reports (CSV/Excel/PDF)",
            "Web retrieval for unit enrichment",
            "Audit trail for all decisions"
        ],
        "Technical Stack",
        [
            "Python 3.11+",
            "Streamlit (UI)",
            "SQLite (Database)",
            "Playwright (Web retrieval)",
            "MCP backend integration",
            "LLM assistance hooks"
        ]
    )
    
    # ========== SLIDE 8: CPL Workflow ==========
    add_content_slide(
        "CPL Automation — Workflow",
        [
            "Step 1: Load SHEA Units",
            "   • Import from Excel (data/SHEA Course Data.xlsx)",
            "",
            "Step 2: Upload Transcript",
            "   • Parse external university transcript PDF",
            "   • Extract unit information",
            "",
            "Step 3: Enrich External Units",
            "   • Crawl external course websites via MCP",
            "   • Gather detailed unit information",
            "",
            "Step 4: Generate Matching Suggestions",
            "   • AI-powered unit matching",
            "   • Confidence scoring with breakdown",
            "",
            "Step 5: Review & Export",
            "   • Approve/reject/override suggestions",
            "   • Export final reports"
        ]
    )
    
    # ========== SLIDE 9: CPL Confidence Scoring ==========
    add_content_slide(
        "CPL Automation — Confidence Scoring",
        [
            "Weighted Components:",
            "",
            "• Title Similarity — Course/unit name matching",
            "• Description Similarity — Content overlap analysis",
            "• Learning Outcomes Similarity — Skill alignment",
            "• Credit Similarity — Credit point comparison",
            "• Grade Bonus — Performance consideration",
            "• Retrieval Bonus — Verified external data",
            "",
            "Quality Controls:",
            "• Component percentages visible for audit",
            "• Non-passing grades flagged (Fail/Not Competent)",
            "• Manual review required for low-confidence matches",
            "• Full decision history tracked in database"
        ]
    )
    
    # ========== SLIDE 10: Project 3 - Title ==========
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    title_shape.text = "Project 3: Student Risk Analytics"
    subtitle_shape = slide.placeholders[1]
    subtitle_shape = "Early Intervention Dashboard"
    
    # ========== SLIDE 11: Dashboard Features ==========
    add_two_column_slide(
        "Risk Analytics — Key Features",
        "Capabilities",
        [
            "Moodle API integration (courses, quizzes, grades)",
            "Student engagement analysis from logs",
            "Risk scoring (0-100 scale)",
            "Auto-categorization (Critical/Warning/Safe)",
            "Direct email outreach via Gmail SMTP",
            "Interactive Streamlit dashboard"
        ],
        "Technical Stack",
        [
            "Python 3.13",
            "Streamlit",
            "Pandas & NumPy",
            "Plotly (visualizations)",
            "Moodle Web Services API",
            "Gmail SMTP for outreach"
        ]
    )
    
    # ========== SLIDE 12: Risk Calculation ==========
    add_content_slide(
        "Risk Analytics — Scoring Methodology",
        [
            "Risk Score Formula (0-100):",
            "",
            "Risk Score = 100 - (0.3×Engagement + 0.4×Completion + 0.3×Performance)",
            "",
            "Components:",
            "• Engagement Score (30%) — Moodle activity logs (clicks + dwell time)",
            "• Assessment Completion (40%) — Submitted quizzes & assignments",
            "• Performance (30%) — Current weighted average grade",
            "",
            "Categorization:",
            "• Critical: Score >75 OR 3+ missed quizzes OR 2+ missed assignments",
            "• Warning: Score 50-75 OR 1+ missed quiz/assignment",
            "• Safe: Score <50",
            "",
            "Perfect student = Risk Score of 0.00"
        ]
    )
    
    # ========== SLIDE 13: Dashboard Intervention ==========
    add_content_slide(
        "Risk Analytics — Intervention Tools",
        [
            "Student Detail View:",
            "• Individual performance breakdown",
            "• Missing submissions visibility",
            "• Quiz attempt history",
            "• Engagement metrics",
            "",
            "Outreach Capabilities:",
            "• Filter students by risk score/category",
            "• Select multiple students",
            "• Send targeted intervention emails",
            "• Gmail SMTP with app password authentication",
            "• Track outreach history",
            "",
            "Privacy & Compliance:",
            "• Follows institutional data privacy policies",
            "• Secure credential management via .env"
        ]
    )
    
    # ========== SLIDE 14: Comparison ==========
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    title_tf = title_shape.text_frame
    title_tf.text = "Project Comparison"
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.size = Pt(32)
    
    # Create table
    rows = 4
    cols = 4
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(12.3)
    height = Inches(2.5)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Set column widths
    table.columns[0].width = Inches(3)
    table.columns[1].width = Inches(3.1)
    table.columns[2].width = Inches(3.1)
    table.columns[3].width = Inches(3.1)
    
    # Headers
    headers = ["Feature", "HATBOT", "CPL Automation", "Risk Analytics"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
    
    # Data
    data = [
        ["Primary Purpose", "Student Support", "Credit Transfer", "Early Intervention"],
        ["AI/ML", "GPT-4o & Gemini", "LLM Assistance", "Risk Scoring Algorithm"],
        ["Data Source", "FAQ Documents", "Transcripts + Web", "Moodle API"]
    ]
    
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_data
            cell.text_frame.paragraphs[0].font.size = Pt(12)
    
    # ========== SLIDE 15: Tech Stack Summary ==========
    add_content_slide(
        "Technology Stack Summary",
        [
            "Frontend Technologies:",
            "• Next.js 16 + React 19 (HATBOT)",
            "• Streamlit (CPL Automation & Risk Analytics)",
            "• Tailwind CSS 4 (HATBOT)",
            "• Plotly (Risk Analytics visualizations)",
            "",
            "Backend & Data:",
            "• Python 3.11-3.13",
            "• TypeScript 5",
            "• MongoDB (HATBOT)",
            "• SQLite (CPL & Risk Analytics)",
            "",
            "Integrations:",
            "• OpenAI GPT-4o-mini & Google Gemini APIs",
            "• Moodle Web Services API",
            "• Gmail SMTP",
            "• MCP Backend (CPL)",
            "• Playwright (Web retrieval)"
        ]
    )
    
    # ========== SLIDE 16: Impact & Benefits ==========
    add_content_slide(
        "Impact & Benefits",
        [
            "Operational Efficiency:",
            "• Automates repetitive administrative tasks",
            "• Reduces manual workload for staff",
            "• 24/7 student support availability",
            "",
            "Student Experience:",
            "• Instant answers to common questions",
            "• Faster credit transfer processing",
            "• Proactive support for at-risk students",
            "",
            "Data-Driven Decisions:",
            "• Confidence scoring with audit trails",
            "• Risk analytics for early intervention",
            "• Exportable reports for compliance",
            "",
            "Scalability:",
            "• All systems designed for growth",
            "• Cloud-ready deployment options",
            "• Modular architecture for extensions"
        ]
    )
    
    # ========== SLIDE 17: Next Steps ==========
    add_content_slide(
        "Next Steps & Future Enhancements",
        [
            "SHEAC-HATBOT:",
            "• Expand FAQ knowledge base",
            "• Add voice interaction support",
            "• Multi-language support",
            "",
            "CPL Automation:",
            "• OAuth-based email integration",
            "• Enhanced LLM reasoning layer",
            "• Batch processing for large transcripts",
            "",
            "Risk Analytics:",
            "• Machine learning risk prediction",
            "• Real-time Moodle sync",
            "• Role-based access control",
            "",
            "Cross-Project:",
            "• Unified authentication system",
            "• Shared analytics dashboard",
            "• Mobile app interfaces"
        ]
    )
    
    # ========== SLIDE 18: Q&A ==========
    add_title_slide(
        "Questions?",
        "Thank you!\n\nContact: Sunil Paudel & Xueliang Lu"
    )
    
    # Save presentation
    output_path = "/Users/alexlu/Desktop/shea-projects/SHEA_Projects_Presentation.pptx"
    prs.save(output_path)
    return output_path

if __name__ == "__main__":
    output = create_presentation()
    print(f"✅ Presentation created: {output}")
